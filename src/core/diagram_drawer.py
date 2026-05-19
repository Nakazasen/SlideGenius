"""Primitive diagram drawing helpers for layout-aware PPTX renderers."""
from typing import Any, Dict, Iterable, List, Tuple

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a hex string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


class DiagramDrawer:
    """Draws primitive nodes and connectors inside a slide."""

    def __init__(self, slide):
        self.slide = slide
        self.shapes = slide.shapes

    def draw_node(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        text: str,
        fill_color: str,
        text_color: str,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        font_size: int = 14,
    ):
        """Draw a text node."""
        shape = self.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        shape.line.color.rgb = hex_to_rgb(fill_color)
        shape.line.width = Pt(1)
        frame = shape.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(text_color)
        p.font.name = "Arial"
        return shape

    def draw_connector(self, shape_from, shape_to, color: str):
        """Draw a connector between two nodes."""
        connector = self.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0)
        connector.begin_connect(shape_from, 1)
        connector.end_connect(shape_to, 3)
        connector.line.color.rgb = hex_to_rgb(color)
        connector.line.width = Pt(2)
        connector.line.end_arrowhead = MSO_LINE.ARROW
        return connector

    def draw_process_flow(
        self,
        steps: Iterable[str],
        area: Tuple[float, float, float, float],
        primary_color: str,
        accent_color: str,
    ) -> None:
        """Draw a simple horizontal process flow within a bounding box."""
        steps = [step for step in steps if step]
        if not steps:
            return

        left, top, width, height = area
        gap = 0.2
        node_w = max(1.5, (width - gap * (len(steps) - 1)) / len(steps))
        node_h = min(1.0, height * 0.55)
        y = top + (height - node_h) / 2
        previous_shape = None

        for index, step in enumerate(steps):
            x = left + index * (node_w + gap)
            fill = accent_color if index == len(steps) - 1 else primary_color
            shape = self.draw_node(x, y, node_w, node_h, step, fill, "#FFFFFF", font_size=12)
            if previous_shape:
                self.draw_connector(previous_shape, shape, "#94A3B8")
            previous_shape = shape

    def draw_comparison(
        self,
        sections: List[Dict[str, Any]],
        area: Tuple[float, float, float, float],
        before_color: str,
        after_color: str,
        neutral_color: str,
    ) -> None:
        """Draw a before/after comparison inside a bounding box."""
        if len(sections) < 2:
            return

        left, top, width, height = area
        col_gap = 0.35
        col_w = (width - col_gap) / 2
        header_h = 0.55
        item_gap = 0.16
        max_items = max(len(sections[0].get("items", [])), len(sections[1].get("items", [])), 1)
        item_h = min(0.72, max(0.45, (height - header_h - item_gap * (max_items + 1)) / max_items))

        def draw_column(section: Dict[str, Any], x: float, color: str):
            self.draw_node(x, top, col_w, header_h, section.get("title", ""), neutral_color, "#FFFFFF", font_size=12)
            y = top + header_h + item_gap
            prev = None
            for item in section.get("items", []):
                shape = self.draw_node(x, y, col_w, item_h, item, color, "#FFFFFF", font_size=11)
                if prev:
                    connector = self.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0)
                    connector.begin_connect(prev, 2)
                    connector.end_connect(shape, 0)
                    connector.line.color.rgb = hex_to_rgb(neutral_color)
                    connector.line.width = Pt(1.2)
                prev = shape
                y += item_h + item_gap

        draw_column(sections[0], left, before_color)
        draw_column(sections[1], left + col_w + col_gap, after_color)
