from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.core.pptx_generator import PPTXGenerator
from src.data.models import Outline, SlideItem, SlideType


def _text_runs(slide):
    runs = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                runs.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        runs.append(text)
    return runs


def _has_text_or_graphic(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            return True
        if getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
            return True
        if shape.shape_type in {
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
            MSO_SHAPE_TYPE.GROUP,
            MSO_SHAPE_TYPE.LINE,
            MSO_SHAPE_TYPE.PICTURE,
        }:
            return True
    return False


def _assert_valid_pptx(output_path: Path, expected_slide_count: int):
    assert output_path.exists()
    assert output_path.stat().st_size > 0
    package = Presentation(str(output_path))
    assert len(package.slides) == expected_slide_count
    for slide in package.slides:
        assert _has_text_or_graphic(slide)
    return package


def _base_slides():
    return [
        SlideItem(
            title="Golden Layout Export",
            slide_type=SlideType.TITLE,
            subtitle="Offline renderer regression deck",
            key_message="Editable PPTX output remains the export contract.",
        ),
        SlideItem(
            title="Section: Context",
            slide_type=SlideType.SECTION,
            summary="Why this export pass matters.",
        ),
        SlideItem(
            title="Content Slide",
            layout_variant="content_2col",
            summary="A compact business update with plain bullets.",
            bullets=["Offline export must be deterministic", "Slides remain editable", "Warnings are persisted on the slide model"],
            key_message="Content slides provide the safe fallback baseline.",
        ),
        SlideItem(
            title="Image Left Alias",
            layout_variant="content_2col",
            summary="Current model has no dedicated image-left renderer, so content locks current fallback behavior.",
            bullets=["No live image generation", "No rasterized slide", "Text remains editable"],
        ),
        SlideItem(
            title="Image Right Fallback",
            layout_variant="content_image_right",
            summary="Missing image assets must fall back to editable text.",
            bullets=["Image unavailable", "Fallback warning expected", "No network call required"],
            image_path="missing-image-for-offline-test.png",
        ),
        SlideItem(
            title="Before and After",
            layout_variant="comparison_before_after",
            summary="Comparison layout uses editable diagram primitives.",
            sections=[
                {"title": "Before", "items": ["Manual", "Slow", "Opaque"]},
                {"title": "After", "items": ["Automated", "Fast", "Visible"]},
            ],
            key_message="The comparison renderer should not crash.",
        ),
        SlideItem(
            title="Process Flow",
            layout_variant="process_flow",
            summary="Process layout uses editable nodes and connectors.",
            sections=[{"title": "Steps", "items": ["Intake", "Assess", "Prioritize", "Execute", "Review"]}],
        ),
        SlideItem(
            title="Timeline Alias",
            layout_variant="process_flow",
            summary="Timeline behavior is currently covered by the process renderer.",
            bullets=["Q1 discover", "Q2 build", "Q3 launch", "Q4 scale"],
        ),
        SlideItem(
            title="KPI Grid",
            layout_variant="kpi_grid",
            stats=[
                {"label": "Revenue", "value": "+18%", "insight": "Growth improved"},
                {"label": "Margin", "value": "+3pt", "insight": "Quality improved"},
                {"label": "Cycle", "value": "-9%", "insight": "Faster delivery"},
                {"label": "NPS", "value": "62", "insight": "Sentiment healthy"},
            ],
        ),
        SlideItem(
            title="Chart Focus",
            layout_variant="chart_focus",
            chart_spec={
                "type": "column",
                "series": [
                    {"label": "Q1", "value": 12},
                    {"label": "Q2", "value": 19},
                    {"label": "Q3", "value": 27},
                    {"label": "Q4", "value": 34},
                ],
                "insight": "Growth accelerates through the year.",
            },
            bullets=["Chart should be a native editable PowerPoint chart."],
        ),
        SlideItem(
            title="Table Summary",
            layout_variant="table_summary",
            table_spec={
                "headers": ["Workstream", "Decision", "Owner"],
                "rows": [
                    ["Sales", "Invest", "CRO"],
                    ["Ops", "Automate", "COO"],
                    ["Support", "Consolidate", "CS"],
                ],
            },
            key_message="Tables should remain native and editable.",
        ),
        SlideItem(
            title="Quote Layout Alias",
            slide_type=SlideType.QUOTE,
            layout_variant="content_2col",
            summary="\"Clarity beats cleverness in export regression tests.\"",
            bullets=["Quote slide currently renders through content fallback."],
        ),
        SlideItem(
            title="Closing Action",
            slide_type=SlideType.CLOSING,
            key_message="Keep renderer behavior stable before larger refactors.",
            bullets=["Approve SG-004 contract tests next"],
        ),
    ]


def test_export_layout_variant_golden_deck_opens_and_has_expected_content(tmp_path: Path):
    outline = Outline(title="Golden Layout Export", slides=_base_slides())
    output_path = tmp_path / "golden_layout_export.pptx"

    PPTXGenerator().generate(outline, output_path)
    package = _assert_valid_pptx(output_path, expected_slide_count=len(outline.slides))

    first_slide_text = "\n".join(_text_runs(package.slides[0]))
    assert "Golden Layout Export" in first_slide_text

    content_text = "\n".join(_text_runs(package.slides[2]))
    assert "Offline export must be deterministic" in content_text
    assert "Slides remain editable" in content_text

    assert any("Image unavailable" in warning for warning in outline.slides[4].render_warnings)
    assert any(shape.has_chart for shape in package.slides[9].shapes if hasattr(shape, "has_chart"))
    assert any(getattr(shape, "has_table", False) for shape in package.slides[10].shapes)


def test_export_unsupported_layout_falls_back_with_warning(tmp_path: Path):
    unsupported = SlideItem(
        title="Unsupported Layout",
        layout_variant="experimental_canvas",
        bullets=["Fallback point one", "Fallback point two"],
        key_message="Unsupported variants must not be silent.",
    )
    outline = Outline(title="Unsupported Layout Deck", slides=[unsupported])
    output_path = tmp_path / "unsupported_layout.pptx"

    PPTXGenerator().generate(outline, output_path)
    package = _assert_valid_pptx(output_path, expected_slide_count=1)

    slide_text = "\n".join(_text_runs(package.slides[0]))
    assert "Fallback point one" in slide_text
    assert any("Unsupported layout variant" in warning for warning in unsupported.render_warnings)


def test_export_dense_content_adds_warnings_and_review_banner(tmp_path: Path):
    long_bullet = "This is a deliberately long bullet that should trigger fitting logic and reduce body size while keeping the exported slide editable."
    dense_slide = SlideItem(
        title="Dense Content",
        layout_variant="content_2col",
        summary="Dense content should not crash export.",
        bullets=[
            long_bullet,
            "Second dense point with enough detail to matter",
            "Third dense point with enough detail to matter",
            "Fourth dense point with enough detail to matter",
            "Fifth point should move into notes and trigger review",
            "Sixth point should also move into notes and trigger review",
        ],
    )
    outline = Outline(title="Dense Content Deck", slides=[dense_slide])
    output_path = tmp_path / "dense_content.pptx"

    PPTXGenerator().generate(outline, output_path)
    package = _assert_valid_pptx(output_path, expected_slide_count=1)

    assert dense_slide.render_warnings
    assert any("Content exceeds layout budget" in warning for warning in dense_slide.render_warnings)
    assert any("Long bullet text reduced" in warning for warning in dense_slide.render_warnings)
    assert "Additional points:" in dense_slide.speaker_notes
    slide_text = "\n".join(_text_runs(package.slides[0]))
    assert "Slide có nhiều ý hơn sức chứa an toàn" in slide_text


def test_export_chart_table_and_diagram_slides_are_native_or_editable(tmp_path: Path):
    outline = Outline(
        title="Native Editable Checks",
        slides=[
            SlideItem(
                title="Native Chart",
                layout_variant="chart_focus",
                chart_spec={
                    "type": "line",
                    "series": [
                        {"label": "Jan", "value": 8},
                        {"label": "Feb", "value": 12},
                        {"label": "Mar", "value": 16},
                    ],
                },
            ),
            SlideItem(
                title="Native Table",
                layout_variant="table_summary",
                table_spec={"headers": ["Area", "Status"], "rows": [["Export", "Green"], ["QA", "Ready"]]},
            ),
            SlideItem(
                title="Editable Diagram",
                layout_variant="process_flow",
                sections=[{"title": "Steps", "items": ["Plan", "Build", "Verify"]}],
            ),
        ],
    )
    output_path = tmp_path / "native_editable_checks.pptx"

    PPTXGenerator().generate(outline, output_path)
    package = _assert_valid_pptx(output_path, expected_slide_count=3)

    assert any(shape.has_chart for shape in package.slides[0].shapes if hasattr(shape, "has_chart"))
    assert any(getattr(shape, "has_table", False) for shape in package.slides[1].shapes)
    diagram_text = "\n".join(_text_runs(package.slides[2]))
    assert "Plan" in diagram_text
    assert "Build" in diagram_text
