"""Theme-aware PPTX generator with semantic slide renderers."""
from __future__ import annotations

from datetime import datetime
from pathlib import Path
import re
from typing import Callable, Dict, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.core.diagram_drawer import DiagramDrawer, hex_to_rgb
from src.core.image_generator import ImageGenerator
from src.core.presentation_spec import FitResult, fit_slide_content, normalize_outline
from src.data.models import Outline, SlideItem, Template


class PPTXGenerator:
    """Generate PowerPoint presentations from semantic slide specs."""

    DEFAULT_TEMPLATE = Template(
        name="executive_blue",
        display_name="Executive Blue",
        description="Structured executive deck with strong blue accents",
        category="Business",
        primary_color="#1D4ED8",
        secondary_color="#0F172A",
        accent_color="#EA580C",
        background_color="#FFFFFF",
        surface_color="#F8FAFC",
        muted_text_color="#64748B",
        font_heading="Arial",
        font_body="Arial",
        title_scale=1.0,
        body_scale=1.0,
        corner_style="rounded",
        accent_style="bar",
    )

    def __init__(self, template: Template = None):
        self.template = template or self.DEFAULT_TEMPLATE
        self.prs: Optional[Presentation] = None
        self.image_generator = None
        self._occupied_regions = []
        self._current_slide_item: Optional[SlideItem] = None
        self._load_template_settings()
        self.renderers: Dict[str, Callable] = {
            "title_hero": self.render_title_hero,
            "agenda": self.render_agenda,
            "section_break": self.render_section_break,
            "content_2col": self.render_content_2col,
            "content_image_right": self.render_content_image_right,
            "stats_highlight": self.render_stats_highlight,
            "kpi_grid": self.render_kpi_grid,
            "chart_focus": self.render_chart_focus,
            "comparison_before_after": self.render_comparison_before_after,
            "decision_matrix": self.render_decision_matrix,
            "process_flow": self.render_process_flow,
            "table_summary": self.render_table_summary,
            "closing_cta": self.render_closing_cta,
        }

    def _load_template_settings(self):
        """Apply user font preferences on top of the selected theme."""
        from src.data.config_manager import ConfigManager

        config = ConfigManager()
        heading_font = config.get("fonts.heading_font", self.template.font_heading)
        body_font = config.get("fonts.body_font", self.template.font_body)
        self.template.font_heading = heading_font
        self.template.font_body = body_font

    def generate(self, outline: Outline, output_path: Path = None) -> Path:
        """Generate PPTX file from outline."""
        outline = normalize_outline(outline)
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        for slide_item in outline.slides:
            normalized_slide, fit = fit_slide_content(slide_item)
            self._add_slide(normalized_slide, fit)

        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = "".join(c for c in outline.title if c.isalnum() or c in " -_")[:30]
            output_path = Path(f"{safe_title}_{timestamp}.pptx")
        else:
            output_path = Path(output_path)

        if output_path.parent and not output_path.parent.exists():
            output_path.parent.mkdir(parents=True, exist_ok=True)

        self.prs.save(str(output_path))
        return output_path

    def _add_slide(self, slide_item: SlideItem, fit: FitResult):
        """Add one slide using the registered renderer with fallback."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._occupied_regions = []
        self._current_slide_item = slide_item
        renderer = self.renderers.get(slide_item.layout_variant)
        if renderer is None:
            slide_item.render_warnings.append(
                f"Unsupported layout variant '{slide_item.layout_variant}'. Fallback to content layout."
            )
            fit.fallback_to_text = True
            renderer = self.render_content_2col
        try:
            renderer(slide, slide_item, fit)
        except Exception:
            slide_item.render_warnings.append(f"Renderer fallback from {slide_item.layout_variant}.")
            fit.fallback_to_text = True
            self.render_content_2col(slide, slide_item, fit)
        self._render_review_banner(slide, fit, slide_item)
        self._add_speaker_notes(slide, slide_item)
        self._current_slide_item = None

    def _add_speaker_notes(self, slide, slide_item: SlideItem):
        """Store notes in slide notes section when possible."""
        notes_text = (slide_item.speaker_notes or slide_item.notes_short or "").strip()
        if not notes_text:
            return
        try:
            text_frame = slide.notes_slide.notes_text_frame
            text_frame.text = notes_text
        except Exception:
            pass

    def _set_background(self, slide, color: str):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(color)

    def _add_textbox(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_size: int,
        color: str,
        font_name: str,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        vertical_anchor=MSO_VERTICAL_ANCHOR.TOP,
    ):
        """Create a simple textbox."""
        estimated_height = self._estimate_text_height(text, font_size, width)
        top, height = self._reserve_safe_region(left, top, width, max(height, estimated_height))
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = vertical_anchor
        frame.margin_top = Pt(2)
        frame.margin_bottom = Pt(2)
        frame.margin_left = Pt(2)
        frame.margin_right = Pt(2)
        p = frame.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = hex_to_rgb(color)
        p.font.name = font_name
        try:
            frame.fit_text(font_family=font_name, max_size=font_size, bold=bold)
        except Exception:
            pass
        return box

    def _add_bullet_list(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        bullets,
        font_size: int,
        color: str,
        font_name: str,
    ):
        """Create a bullet list."""
        estimated_height = self._estimate_bullet_height(bullets, font_size, width)
        top, height = self._reserve_safe_region(left, top, width, max(height, estimated_height))
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_top = Pt(2)
        frame.margin_bottom = Pt(2)
        frame.margin_left = Pt(2)
        frame.margin_right = Pt(2)
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = f"- {bullet}"
            paragraph.level = 0
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = hex_to_rgb(color)
            paragraph.font.name = font_name
            paragraph.space_after = Pt(10)
        try:
            frame.fit_text(font_family=font_name, max_size=font_size)
        except Exception:
            pass
        return box

    def _estimate_text_height(self, text: str, font_size: int, width: float) -> float:
        """Roughly estimate textbox height in inches."""
        content = str(text or "").strip()
        if not content:
            return 0.3
        chars_per_line = max(10, int(width * 8.5))
        line_count = 0
        for raw_line in content.splitlines() or [content]:
            line = raw_line.strip() or " "
            line_count += max(1, (len(line) // chars_per_line) + 1)
        return max(0.32, line_count * max(0.24, font_size / 42.0))

    def _estimate_bullet_height(self, bullets, font_size: int, width: float) -> float:
        """Estimate bullet list height in inches."""
        if not bullets:
            return 0.35
        total = 0.1
        for bullet in bullets:
            total += self._estimate_text_height(f"- {bullet}", font_size, width) + 0.05
        return total

    def _reserve_safe_region(self, left: float, top: float, width: float, height: float) -> tuple[float, float]:
        """Shift text regions downward to reduce overlap and clamp them within the slide."""
        slide_bottom = 7.15
        gap = 0.08
        adjusted_top = top
        for _ in range(12):
            overlaps = [
                region for region in self._occupied_regions
                if self._regions_overlap((left, adjusted_top, width, height), region)
            ]
            if not overlaps:
                break
            adjusted_top = max(region[3] for region in overlaps) + gap
        remaining = max(0.35, slide_bottom - adjusted_top)
        adjusted_height = min(height, remaining)
        if adjusted_height < height and self._current_slide_item is not None:
            warning = "Textbox height had to be clamped to avoid overflow. Manual review recommended."
            if warning not in self._current_slide_item.render_warnings:
                self._current_slide_item.render_warnings.append(warning)
        region = (left, adjusted_top, left + width, adjusted_top + adjusted_height)
        self._occupied_regions.append(region)
        return adjusted_top, adjusted_height

    def _regions_overlap(self, first: tuple[float, float, float, float], second: tuple[float, float, float, float]) -> bool:
        """Return True when two slide regions intersect."""
        return not (
            first[2] <= second[0]
            or first[0] >= second[2]
            or first[3] <= second[1]
            or first[1] >= second[3]
        )

    def _fit_bullets(self, slide_item: SlideItem, fit: FitResult, limit: Optional[int] = None):
        """Return the safe bullet subset for this slide render."""
        bullets = fit.visible_bullets or slide_item.normalized_bullets
        return bullets[:limit] if limit is not None else bullets

    def _fit_stats(self, slide_item: SlideItem, fit: FitResult, limit: Optional[int] = None):
        """Return the safe stats subset for this slide render."""
        stats = fit.visible_stats or slide_item.stats
        return stats[:limit] if limit is not None else stats

    def _fit_sections(self, slide_item: SlideItem, fit: FitResult, limit: Optional[int] = None):
        """Return the safe sections subset for this slide render."""
        sections = fit.visible_sections or slide_item.sections
        return sections[:limit] if limit is not None else sections

    def _fit_table_rows(self, slide_item: SlideItem, fit: FitResult):
        """Return the safe table row subset for this slide render."""
        if fit.visible_table_rows:
            return fit.visible_table_rows
        return slide_item.table_spec.get("rows", []) if slide_item.table_spec else []

    def _render_review_banner(self, slide, fit: FitResult, slide_item: SlideItem):
        """Show an explicit on-slide review banner when content exceeded the safe layout budget."""
        if not fit.requires_manual_review:
            return
        message = fit.overflow_notice or "Slide cần rà soát thủ công vì nội dung vượt ngân sách bố cục."
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(6.92), Inches(12.15), Inches(0.28))
        banner.fill.solid()
        banner.fill.fore_color.rgb = hex_to_rgb("#FEF3C7")
        banner.line.fill.background()
        box = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(11.85), Inches(0.2))
        frame = box.text_frame
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.text = message
        paragraph.font.size = Pt(10)
        paragraph.font.bold = True
        paragraph.font.color.rgb = hex_to_rgb("#92400E")
        paragraph.font.name = self.template.font_body

    def _add_surface(self, slide, left: float, top: float, width: float, height: float, color: str):
        """Draw a rounded surface card."""
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(color)
        shape.line.color.rgb = hex_to_rgb(color)
        return shape

    def _accent_bar(self, slide, left: float, top: float, width: float, color: str, height: float = 0.08):
        """Draw a thin accent rule."""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(color)
        bar.line.fill.background()
        return bar

    def _set_cell_text(self, cell, text: str, font_size: int, color: str, bold: bool = False, align=PP_ALIGN.LEFT):
        """Apply consistent table cell styling."""
        cell.text = str(text)
        frame = cell.text_frame
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = hex_to_rgb(color)
        paragraph.font.bold = bold
        paragraph.font.name = self.template.font_body

    def _semantic_color(self, value: str = "", insight: str = "") -> str:
        """Choose a semantic accent based on positive or negative signal."""
        text = f"{value} {insight}".lower()
        positive_markers = ["+", "up", "improve", "growth", "faster", "higher", "healthy", "better"]
        negative_markers = ["-", "down", "risk", "decline", "slower", "lower", "pressure", "worse"]
        if any(marker in text for marker in negative_markers):
            return "#DC2626"
        if any(marker in text for marker in positive_markers):
            return "#059669"
        return self.template.primary_color

    def _apply_micro_polish(self, slide_item: SlideItem):
        """Mark layouts that received exporter-side visual refinement."""
        slide_item.micro_polish_applied = slide_item.layout_variant in {"stats_highlight", "kpi_grid", "chart_focus", "table_summary"}

    def _resolve_image(self, slide_item: SlideItem) -> Optional[str]:
        """Return an existing or generated image path."""
        if slide_item.image_path and Path(slide_item.image_path).exists():
            return slide_item.image_path
        if not slide_item.image_prompt:
            return None
        if self.image_generator is None:
            self.image_generator = ImageGenerator()
        try:
            result = self.image_generator.generate_image(slide_item.image_prompt)
            return result if result and Path(result).exists() else None
        except Exception:
            return None

    def _to_numeric(self, value) -> float:
        """Coerce mixed display values into numbers for chart rendering."""
        if isinstance(value, (int, float)):
            return float(value)
        text = str(value or "").strip()
        if not text:
            return 0.0
        cleaned = re.sub(r"[^0-9.\-]+", "", text.replace(",", ""))
        try:
            return float(cleaned)
        except ValueError:
            return 0.0

    def _normalize_chart_spec(self, slide_item: SlideItem, chart_spec_override: Optional[dict] = None) -> Optional[dict]:
        """Normalize chart payload into categories and numeric series sets."""
        chart_spec = chart_spec_override or slide_item.chart_spec or {}
        chart_type = str(chart_spec.get("type", "bar")).strip().lower()
        series_groups = []

        raw_groups = chart_spec.get("series_groups", [])
        if isinstance(raw_groups, list) and raw_groups:
            for index, group in enumerate(raw_groups[:2]):
                if not isinstance(group, dict):
                    continue
                points = []
                for point in group.get("points", [])[:6]:
                    if isinstance(point, dict):
                        points.append(
                            {
                                "label": str(point.get("label", "")).strip(),
                                "value": self._to_numeric(point.get("value", 0)),
                            }
                        )
                if points:
                    series_groups.append(
                        {
                            "name": str(group.get("name", f"Series {index + 1}")).strip() or f"Series {index + 1}",
                            "points": points,
                        }
                    )
        else:
            points = []
            for point in chart_spec.get("series", [])[:6]:
                if isinstance(point, dict):
                    points.append(
                        {
                            "label": str(point.get("label", "")).strip(),
                            "value": self._to_numeric(point.get("value", 0)),
                        }
                    )
            if not points and slide_item.stats:
                points = [
                    {
                        "label": str(stat.get("label", f"Metric {index + 1}")).strip(),
                        "value": self._to_numeric(stat.get("value", index + 1)),
                    }
                    for index, stat in enumerate(slide_item.stats[:4])
                ]
            if points:
                series_groups.append({"name": str(chart_spec.get("series_name", "Series")).strip() or "Series", "points": points})

        if not series_groups:
            return None

        categories = [point["label"] or f"Item {index + 1}" for index, point in enumerate(series_groups[0]["points"])]
        values = []
        for group in series_groups:
            row = []
            for cat_index, category in enumerate(categories):
                match = next((point for point in group["points"] if point["label"] == category), None)
                row.append(match["value"] if match else 0.0)
            values.append({"name": group["name"], "values": row})

        return {
            "type": chart_type,
            "categories": categories,
            "series_sets": values,
            "insight": str(chart_spec.get("insight", "")).strip(),
        }

    def _add_native_chart(self, slide, slide_item: SlideItem, fit: FitResult, left: float, top: float, width: float, height: float):
        """Render a native PowerPoint chart from chart_spec."""
        payload = self._normalize_chart_spec(slide_item, fit.visible_chart_spec)
        if not payload:
            return None

        chart_type_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "donut": XL_CHART_TYPE.DOUGHNUT,
        }
        chart_type = chart_type_map.get(payload["type"], XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart_data = CategoryChartData()
        chart_data.categories = payload["categories"]
        for series_set in payload["series_sets"]:
            chart_data.add_series(series_set["name"], series_set["values"])

        graphic_frame = slide.shapes.add_chart(
            chart_type,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            chart_data,
        )
        chart = graphic_frame.chart
        chart.has_title = False
        chart.has_legend = len(payload["series_sets"]) > 1 or payload["type"] == "donut"
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM if payload["type"] != "donut" else XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False

        palette = [
            self.template.primary_color,
            self.template.accent_color,
            self.template.secondary_color,
            "#0EA5E9",
            "#14B8A6",
            "#F59E0B",
        ]

        try:
            plot = chart.plots[0]
            plot.has_data_labels = payload["type"] in {"bar", "column", "donut"}
            if payload["type"] in {"bar", "column"}:
                chart.category_axis.tick_labels.font.size = Pt(11)
                chart.value_axis.tick_labels.font.size = Pt(10)
                chart.value_axis.has_major_gridlines = False
                chart.value_axis.minimum_scale = 0
            elif payload["type"] == "line":
                chart.category_axis.tick_labels.font.size = Pt(11)
                chart.value_axis.tick_labels.font.size = Pt(10)
                chart.value_axis.has_major_gridlines = True
            for series_index, series in enumerate(chart.series):
                color = palette[series_index % len(palette)]
                if payload["type"] == "line":
                    series.format.line.color.rgb = hex_to_rgb(color)
                elif payload["type"] == "donut":
                    for point_index, point in enumerate(series.points):
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = hex_to_rgb(palette[point_index % len(palette)])
                else:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = hex_to_rgb(color)
                    series.format.line.color.rgb = hex_to_rgb(color)
        except Exception:
            return chart

        return chart

    def render_title_hero(self, slide, slide_item: SlideItem, fit: FitResult):
        self._set_background(slide, self.template.background_color)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.22))
        accent.fill.solid()
        accent.fill.fore_color.rgb = hex_to_rgb(self.template.primary_color)
        accent.line.fill.background()
        self._add_textbox(
            slide,
            0.9,
            1.45,
            8.9,
            1.6,
            slide_item.title,
            int(fit.title_size * 1.45 * self.template.title_scale),
            self.template.secondary_color,
            self.template.font_heading,
            bold=True,
        )
        subtitle = slide_item.subtitle or slide_item.summary or slide_item.key_message
        if subtitle:
            self._add_textbox(
                slide,
                0.95,
                3.35,
                7.9,
                1.1,
                subtitle,
                int(fit.subtitle_size * self.template.body_scale),
                self.template.muted_text_color,
                self.template.font_body,
            )
        self._add_surface(slide, 9.6, 1.2, 2.7, 4.8, self.template.surface_color)
        self._add_textbox(
            slide,
            9.95,
            2.15,
            2.0,
            1.6,
            slide_item.key_message or slide_item.summary or "Executive summary",
            24,
            self.template.primary_color,
            self.template.font_heading,
            bold=True,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        )

    def render_agenda(self, slide, slide_item: SlideItem, fit: FitResult):
        self._set_background(slide, self.template.background_color)
        self._add_textbox(slide, 0.8, 0.6, 5.2, 0.8, slide_item.title, fit.title_size + 4, self.template.secondary_color, self.template.font_heading, bold=True)
        if slide_item.summary:
            self._add_textbox(slide, 0.8, 1.35, 6.2, 0.6, slide_item.summary, fit.subtitle_size, self.template.muted_text_color, self.template.font_body)
        base_top = 2.0
        for index, bullet in enumerate(self._fit_bullets(slide_item, fit, 6)):
            y = base_top + index * 0.76
            chip = self._add_surface(slide, 0.9, y, 0.7, 0.46, self.template.primary_color)
            chip.text_frame.paragraphs[0].text = f"{index + 1:02d}"
            chip.text_frame.paragraphs[0].font.size = Pt(14)
            chip.text_frame.paragraphs[0].font.bold = True
            chip.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#FFFFFF")
            chip.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._add_textbox(slide, 1.9, y - 0.02, 9.6, 0.56, bullet, fit.body_size, self.template.secondary_color, self.template.font_body)

    def render_section_break(self, slide, slide_item: SlideItem, fit: FitResult):
        self._set_background(slide, self.template.secondary_color)
        self._add_textbox(slide, 0.9, 1.6, 11.6, 1.0, slide_item.title, fit.title_size + 6, "#FFFFFF", self.template.font_heading, bold=True)
        message = slide_item.key_message or slide_item.summary or slide_item.subtitle
        if message:
            self._add_textbox(slide, 0.95, 3.0, 8.4, 1.1, message, fit.subtitle_size + 2, "#E2E8F0", self.template.font_body)
        self._add_surface(slide, 10.2, 1.2, 1.1, 4.8, self.template.accent_color)

    def render_content_2col(self, slide, slide_item: SlideItem, fit: FitResult):
        self._set_background(slide, self.template.background_color)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.55), Inches(0.16), Inches(0.9))
        accent.fill.solid()
        accent.fill.fore_color.rgb = hex_to_rgb(self.template.primary_color)
        accent.line.fill.background()
        self._add_textbox(slide, 1.1, 0.55, 10.8, 0.85, slide_item.title, fit.title_size + 2, self.template.secondary_color, self.template.font_heading, bold=True)
        if slide_item.summary or slide_item.subtitle:
            self._add_textbox(slide, 1.1, 1.35, 10.2, 0.55, slide_item.summary or slide_item.subtitle, fit.subtitle_size, self.template.muted_text_color, self.template.font_body)

        bullets = self._fit_bullets(slide_item, fit, 4)
        left_bullets = bullets[:2]
        right_bullets = bullets[2:]
        self._add_bullet_list(slide, 1.0, 2.2, 5.0, 3.8, left_bullets, fit.body_size, self.template.secondary_color, self.template.font_body)
        if right_bullets:
            self._add_bullet_list(slide, 6.5, 2.2, 5.0, 3.8, right_bullets, fit.body_size, self.template.secondary_color, self.template.font_body)
        if slide_item.key_message:
            card = self._add_surface(slide, 0.95, 6.0, 11.2, 0.85, self.template.surface_color)
            self._add_textbox(slide, 1.25, 6.18, 10.5, 0.42, slide_item.key_message, 16, self.template.primary_color, self.template.font_body, bold=True)

    def render_content_image_right(self, slide, slide_item: SlideItem, fit: FitResult):
        image_path = self._resolve_image(slide_item)
        if not image_path:
            fit.fallback_to_text = True
            slide_item.render_warnings.append("Image unavailable. Fallback to text layout.")
            self.render_content_2col(slide, slide_item, fit)
            return
        self._set_background(slide, self.template.background_color)
        self._add_textbox(slide, 0.9, 0.55, 6.0, 0.85, slide_item.title, fit.title_size + 2, self.template.secondary_color, self.template.font_heading, bold=True)
        if slide_item.summary or slide_item.subtitle:
            self._add_textbox(slide, 0.9, 1.3, 6.0, 0.65, slide_item.summary or slide_item.subtitle, fit.subtitle_size, self.template.muted_text_color, self.template.font_body)
        self._add_bullet_list(slide, 0.95, 2.0, 5.6, 4.4, self._fit_bullets(slide_item, fit, 4), fit.body_size, self.template.secondary_color, self.template.font_body)
        slide.shapes.add_picture(str(image_path), Inches(7.2), Inches(1.35), width=Inches(5.1), height=Inches(4.3))
        if slide_item.key_message:
            self._add_textbox(slide, 7.25, 5.95, 4.9, 0.6, slide_item.key_message, 14, self.template.primary_color, self.template.font_body, bold=True, align=PP_ALIGN.CENTER)

    def render_stats_highlight(self, slide, slide_item: SlideItem, fit: FitResult):
        self._set_background(slide, self.template.background_color)
        self._apply_micro_polish(slide_item)
        self._add_textbox(slide, 0.8, 0.55, 7.8, 0.8, slide_item.title, fit.title_size + 2, self.template.secondary_color, self.template.font_heading, bold=True)
        if slide_item.summary or slide_item.subtitle:
            self._add_textbox(slide, 0.8, 1.25, 8.0, 0.55, slide_item.summary or slide_item.subtitle, fit.subtitle_size, self.template.muted_text_color, self.template.font_body)
        stats = self._fit_stats(slide_item, fit, 3)
        card_width = 3.8 if len(stats) == 3 else 5.7
        start_x = 0.8
        for index, stat in enumerate(stats):
            x = start_x + index * (card_width + 0.35)
            self._add_surface(slide, x, 2.2, card_width, 2.45, self.template.surface_color)
            metric_color = self._semantic_color(stat.get("value", ""), stat.get("insight", ""))
            self._accent_bar(slide, x, 2.2, card_width, metric_color, 0.07)
            self._add_textbox(slide, x + 0.22, 2.45, card_width - 0.44, 0.4, stat.get("label", "Metric"), 14, self.template.muted_text_color, self.template.font_body)
            self._add_textbox(slide, x + 0.22, 2.92, card_width - 0.44, 0.85, stat.get("value", ""), fit.stat_value_size, metric_color, self.template.font_heading, bold=True)
            self._add_textbox(slide, x + 0.22, 3.8, card_width - 0.44, 0.55, stat.get("insight", ""), 12, self.template.secondary_color, self.template.font_body)
        if slide_item.key_message:
            self._add_surface(slide, 0.8, 5.2, 11.1, 0.72, self.template.surface_color)
            self._add_textbox(slide, 1.0, 5.42, 10.7, 0.34, slide_item.key_message, 15, self.template.secondary_color, self.template.font_body, bold=True)

    def render_kpi_grid(self, slide, slide_item: SlideItem, fit: FitResult):
        self._set_background(slide, self.template.background_color)
        self._apply_micro_polish(slide_item)
        self._add_textbox(slide, 0.8, 0.55, 8.5, 0.8, slide_item.title, fit.title_size + 2, self.template.secondary_color, self.template.font_heading, bold=True)
        if slide_item.summary or slide_item.subtitle:
            self._add_textbox(slide, 0.8, 1.25, 9.0, 0.55, slide_item.summary or slide_item.subtitle, fit.subtitle_size, self.template.muted_text_color, self.template.font_body)
        stats = self._fit_stats(slide_item, fit, 4)
        positions = [(0.8, 2.1), (6.65, 2.1), (0.8, 4.4), (6.65, 4.4)]
        for index, stat in enumerate(stats):
            x, y = positions[index]
            self._add_surface(slide, x, y, 5.15, 1.85, self.template.surface_color)
            metric_color = self._semantic_color(stat.get("value", ""), stat.get("insight", ""))
            self._accent_bar(slide, x, y, 5.15, metric_color, 0.06)
            self._add_textbox(slide, x + 0.24, y + 0.18, 4.65, 0.35, stat.get("label", "KPI"), 13, self.template.muted_text_color, self.template.font_body)
            self._add_textbox(slide, x + 0.24, y + 0.55, 3.0, 0.6, stat.get("value", ""), 24, metric_color, self.template.font_heading, bold=True)
            self._add_textbox(slide, x + 0.24, y + 1.18, 4.65, 0.42, stat.get("insight", ""), 11, self.template.secondary_color, self.template.font_body)
        if slide_item.key_message:
            self._add_surface(slide, 0.8, 6.28, 11.15, 0.5, "#EEF2FF")
            self._add_textbox(slide, 1.0, 6.42, 10.7, 0.24, slide_item.key_message, 14, self.template.accent_color, self.template.font_body, bold=True, align=PP_ALIGN.CENTER)

    def render_chart_focus(self, slide, slide_item: SlideItem, fit: FitResult):
        payload = self._normalize_chart_spec(slide_item)
        if not payload:
            slide_item.render_warnings.append("Chart data missing. Fallback to KPI or content layout.")
            if slide_item.stats:
                self.render_kpi_grid(slide, slide_item, fit)
            else:
                self.render_content_2col(slide, slide_item, fit)
            return
        self._set_background(slide, self.template.background_color)
        self._apply_micro_polish(slide_item)
        self._add_textbox(slide, 0.8, 0.5, 10.4, 0.8, slide_item.title, fit.title_size + 1, self.template.secondary_color, self.template.font_heading, bold=True)
        if slide_item.summary or slide_item.subtitle:
            self._add_textbox(slide, 0.8, 1.18, 9.4, 0.48, slide_item.summary or slide_item.subtitle, fit.subtitle_size, self.template.muted_text_color, self.template.font_body)
        self._add_surface(slide, 0.82, 1.82, 7.45, 4.32, "#FBFDFF")
        self._accent_bar(slide, 0.82, 1.82, 7.45, self.template.primary_color, 0.07)
        chart = self._add_native_chart(slide, slide_item, fit, 0.9, 1.95, 7.2, 3.95)
        if chart is None:
            slide_item.render_warnings.append("Native chart rendering failed. Fallback to content layout.")
            self.render_content_2col(slide, slide_item, fit)
            return
        self._add_surface(slide, 8.45, 1.95, 3.35, 3.95, self.template.surface_color)
        self._accent_bar(slide, 8.45, 1.95, 3.35, self.template.accent_color, 0.07)
        insight = slide_item.key_message or payload.get("insight") or slide_item.summary or "Use the chart to anchor one decision-ready takeaway."
        chart_label = {
            "line": "Trend",
            "bar": "Comparison",
            "column": "Comparison",
            "donut": "Mix",
        }.get(payload["type"], "Chart")
        self._add_textbox(slide, 1.05, 1.88, 1.7, 0.22, chart_label, 11, self.template.primary_color, self.template.font_body, bold=True)
        self._add_textbox(slide, 8.72, 2.2, 2.8, 0.34, "Insight", 12, self.template.muted_text_color, self.template.font_body, bold=True)
        self._add_textbox(slide, 8.72, 2.62, 2.7, 1.08, insight, 16, self.template.primary_color, self.template.font_heading, bold=True)
        support_bullets = self._fit_bullets(slide_item, fit, 2)
        if support_bullets:
            self._add_bullet_list(slide, 8.6, 3.85, 2.95, 1.45, support_bullets, 12, self.template.secondary_color, self.template.font_body)
        if fit.visible_chart_spec and fit.visible_chart_spec.get("type", "").lower() == "donut":
            self._add_textbox(slide, 0.95, 6.15, 6.9, 0.35, "Composition view: read category mix before drilling into actions.", 12, self.template.muted_text_color, self.template.font_body)
        if slide_item.key_message:
            self._add_surface(slide, 0.9, 6.32, 11.0, 0.36, "#FFF7ED")
            self._add_textbox(slide, 1.0, 6.42, 10.8, 0.18, slide_item.key_message, 13, self.template.accent_color, self.template.font_body, bold=True, align=PP_ALIGN.CENTER)

    def render_comparison_before_after(self, slide, slide_item: SlideItem, fit: FitResult):
        sections = self._fit_sections(slide_item, fit, 2)
        if len(sections) < 2:
            slide_item.render_warnings.append("Comparison data incomplete. Fallback to content layout.")
            self.render_content_2col(slide, slide_item, fit)
            return
        self._set_background(slide, self.template.background_color)
        self._add_textbox(slide, 0.8, 0.55, 10.8, 0.8, slide_item.title, fit.title_size + 1, self.template.secondary_color, self.template.font_heading, bold=True)
        if slide_item.summary:
            self._add_textbox(slide, 0.8, 1.28, 11.0, 0.48, slide_item.summary, fit.subtitle_size, self.template.muted_text_color, self.template.font_body)
        drawer = DiagramDrawer(slide)
        drawer.draw_comparison(
            sections,
            area=(0.85, 2.0, 11.6, 3.9),
            before_color="#94A3B8",
            after_color=self.template.primary_color,
            neutral_color=self.template.secondary_color,
        )
        if slide_item.key_message:
            self._add_textbox(slide, 0.95, 6.25, 11.0, 0.45, slide_item.key_message, 14, self.template.accent_color, self.template.font_body, bold=True, align=PP_ALIGN.CENTER)

    def render_decision_matrix(self, slide, slide_item: SlideItem, fit: FitResult):
        sections = self._fit_sections(slide_item, fit, 4)
        if len(sections) < 4:
            slide_item.render_warnings.append("Decision matrix data incomplete. Fallback to comparison layout.")
            self.render_comparison_before_after(slide, slide_item, fit)
            return
        self._set_background(slide, self.template.background_color)
        self._add_textbox(slide, 0.8, 0.45, 10.5, 0.8, slide_item.title, fit.title_size + 1, self.template.secondary_color, self.template.font_heading, bold=True)
        if slide_item.summary:
            self._add_textbox(slide, 0.8, 1.15, 10.8, 0.45, slide_item.summary, fit.subtitle_size, self.template.muted_text_color, self.template.font_body)
        # Axes
        axis_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.25), Inches(1.95), Inches(0.03), Inches(3.8))
        axis_h.fill.solid()
        axis_h.fill.fore_color.rgb = hex_to_rgb(self.template.secondary_color)
        axis_h.line.fill.background()
        axis_v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.15), Inches(3.85), Inches(10.3), Inches(0.03))
        axis_v.fill.solid()
        axis_v.fill.fore_color.rgb = hex_to_rgb(self.template.secondary_color)
        axis_v.line.fill.background()
        quadrants = [(1.2, 2.0), (6.35, 2.0), (1.2, 3.95), (6.35, 3.95)]
        colors = [self.template.surface_color, "#E8F2FF", "#FFF0E8", "#ECFDF5"]
        for index, section in enumerate(sections[:4]):
            x, y = quadrants[index]
            self._add_surface(slide, x, y, 4.85, 1.65, colors[index])
            self._add_textbox(slide, x + 0.18, y + 0.15, 4.45, 0.28, section.get("title", ""), 12, self.template.secondary_color, self.template.font_heading, bold=True)
            items = section.get("items", [])
            for item_index, item in enumerate(items[:2]):
                self._add_textbox(slide, x + 0.18, y + 0.5 + item_index * 0.42, 4.45, 0.28, f"• {item}", 11, self.template.secondary_color, self.template.font_body)
        if slide_item.key_message:
            self._add_textbox(slide, 0.9, 6.35, 11.0, 0.4, slide_item.key_message, 14, self.template.accent_color, self.template.font_body, bold=True, align=PP_ALIGN.CENTER)

    def render_process_flow(self, slide, slide_item: SlideItem, fit: FitResult):
        sections = self._fit_sections(slide_item, fit, 1)
        steps = sections[0].get("items", []) if sections else self._fit_bullets(slide_item, fit, 5)
        if not steps:
            slide_item.render_warnings.append("Process data incomplete. Fallback to content layout.")
            self.render_content_2col(slide, slide_item, fit)
            return
        self._set_background(slide, self.template.background_color)
        self._add_textbox(slide, 0.8, 0.55, 10.8, 0.8, slide_item.title, fit.title_size + 1, self.template.secondary_color, self.template.font_heading, bold=True)
        if slide_item.summary:
            self._add_textbox(slide, 0.8, 1.25, 11.0, 0.5, slide_item.summary, fit.subtitle_size, self.template.muted_text_color, self.template.font_body)
        drawer = DiagramDrawer(slide)
        drawer.draw_process_flow(steps, (0.95, 2.1, 11.2, 2.8), self.template.primary_color, self.template.accent_color)
        if slide_item.key_message:
            self._add_textbox(slide, 1.0, 5.7, 10.8, 0.55, slide_item.key_message, 14, self.template.secondary_color, self.template.font_body, bold=True, align=PP_ALIGN.CENTER)

    def render_table_summary(self, slide, slide_item: SlideItem, fit: FitResult):
        if not slide_item.table_spec:
            slide_item.render_warnings.append("Table data missing. Fallback to content layout.")
            self.render_content_2col(slide, slide_item, fit)
            return
        self._set_background(slide, self.template.background_color)
        self._apply_micro_polish(slide_item)
        self._add_textbox(slide, 0.8, 0.5, 10.8, 0.8, slide_item.title, fit.title_size + 1, self.template.secondary_color, self.template.font_heading, bold=True)
        if slide_item.summary:
            self._add_textbox(slide, 0.8, 1.15, 10.8, 0.45, slide_item.summary, fit.subtitle_size, self.template.muted_text_color, self.template.font_body)
        headers = slide_item.table_spec.get("headers", ["Column 1", "Column 2"])
        rows = self._fit_table_rows(slide_item, fit)
        row_count = max(1, len(rows)) + 1
        col_count = max(2, len(headers))
        self._add_surface(slide, 0.82, 1.85, 11.56, 4.3, "#FCFEFF")
        self._accent_bar(slide, 0.82, 1.85, 11.56, self.template.primary_color, 0.07)
        table_shape = slide.shapes.add_table(row_count, col_count, Inches(0.9), Inches(2.0), Inches(11.4), Inches(3.8))
        table = table_shape.table
        preferred_widths = [2.3, 2.0, 4.9, 1.8]
        for col_index in range(col_count):
            table.columns[col_index].width = Inches(preferred_widths[col_index] if col_index < len(preferred_widths) else 2.0)
        for col_index, header in enumerate(headers[:col_count]):
            cell = table.cell(0, col_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(self.template.primary_color)
            self._set_cell_text(cell, header, 12, "#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        for row_index, row in enumerate(rows[: row_count - 1], start=1):
            for col_index in range(col_count):
                value = row[col_index] if col_index < len(row) else ""
                cell = table.cell(row_index, col_index)
                cell.fill.solid()
                fill_color = self.template.surface_color if row_index % 2 == 1 else "#FFFFFF"
                if col_index == 0:
                    fill_color = "#F1F5F9"
                cell.fill.fore_color.rgb = hex_to_rgb(fill_color)
                self._set_cell_text(
                    cell,
                    str(value),
                    12,
                    self.template.secondary_color if col_index else self.template.primary_color,
                    bold=(col_index == 0),
                )
        if slide_item.key_message:
            self._add_surface(slide, 0.9, 6.1, 11.0, 0.42, "#EEF2FF")
            self._add_textbox(slide, 1.0, 6.22, 10.8, 0.2, slide_item.key_message, 13, self.template.accent_color, self.template.font_body, bold=True, align=PP_ALIGN.CENTER)

    def render_closing_cta(self, slide, slide_item: SlideItem, fit: FitResult):
        self._set_background(slide, self.template.secondary_color)
        self._add_textbox(
            slide,
            1.0,
            1.7,
            11.1,
            1.2,
            slide_item.title,
            fit.title_size + 8,
            "#FFFFFF",
            self.template.font_heading,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        message = slide_item.key_message or slide_item.summary
        if message:
            self._add_textbox(slide, 2.0, 3.35, 9.3, 0.7, message, fit.subtitle_size + 2, "#E2E8F0", self.template.font_body, align=PP_ALIGN.CENTER)
        cta = self._fit_bullets(slide_item, fit, 1)
        if cta:
            pill = self._add_surface(slide, 4.35, 4.8, 4.5, 0.78, self.template.accent_color)
            pill.text_frame.paragraphs[0].text = cta[0]
            pill.text_frame.paragraphs[0].font.size = Pt(18)
            pill.text_frame.paragraphs[0].font.bold = True
            pill.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#FFFFFF")
            pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


TEMPLATES = {
    "executive_blue": Template(
        name="executive_blue",
        display_name="Executive Blue",
        description="Structured executive deck with strong blue accents",
        category="Business",
        primary_color="#1D4ED8",
        secondary_color="#0F172A",
        accent_color="#EA580C",
        background_color="#FFFFFF",
        surface_color="#F8FAFC",
        muted_text_color="#64748B",
        title_scale=1.0,
        body_scale=1.0,
        corner_style="rounded",
        accent_style="bar",
    ),
    "consulting_white": Template(
        name="consulting_white",
        display_name="Consulting White",
        description="Minimal white consulting style deck",
        category="Business",
        primary_color="#0F766E",
        secondary_color="#111827",
        accent_color="#B45309",
        background_color="#FFFFFF",
        surface_color="#F9FAFB",
        muted_text_color="#6B7280",
        title_scale=1.0,
        body_scale=0.98,
        corner_style="soft",
        accent_style="rule",
    ),
    "tech_slate": Template(
        name="tech_slate",
        display_name="Tech Slate",
        description="Dark-leaning tech presentation theme",
        category="Technology",
        primary_color="#2563EB",
        secondary_color="#0F172A",
        accent_color="#F97316",
        background_color="#F8FAFC",
        surface_color="#E2E8F0",
        muted_text_color="#475569",
        title_scale=1.0,
        body_scale=1.0,
        corner_style="rounded",
        accent_style="block",
    ),
    "growth_orange": Template(
        name="growth_orange",
        display_name="Growth Orange",
        description="Growth and GTM theme with warm accents",
        category="Business",
        primary_color="#EA580C",
        secondary_color="#1F2937",
        accent_color="#0EA5E9",
        background_color="#FFFDF8",
        surface_color="#FFEDD5",
        muted_text_color="#7C5E3A",
        title_scale=1.0,
        body_scale=1.0,
        corner_style="rounded",
        accent_style="bar",
    ),
}


TEMPLATES.update(
    {
        "modern_blue": TEMPLATES["executive_blue"],
        "minimal_gray": TEMPLATES["consulting_white"],
        "creative_orange": TEMPLATES["growth_orange"],
        "nature_green": TEMPLATES["consulting_white"],
        "purple_gradient": TEMPLATES["tech_slate"],
    }
)
